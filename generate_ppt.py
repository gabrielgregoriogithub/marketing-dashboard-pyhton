"""
Deel Marketing Analytics Challenge — PPT Generator
Generates a comprehensive executive presentation answering Tasks 1, 2 and 3.
Run: python generate_ppt.py
"""

import io
import os
import numpy as np
import pandas as pd
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
import matplotlib.patches as mpatches
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Brand Colors ──────────────────────────────────────────────────────────────
ORANGE  = RGBColor(0xFF, 0x5C, 0x35)
DARK    = RGBColor(0x14, 0x14, 0x23)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY   = RGBColor(0xF5, 0xF5, 0xF7)
LGRAY2  = RGBColor(0xE8, 0xE8, 0xF0)
GREEN   = RGBColor(0x00, 0xC8, 0x53)
RED_C   = RGBColor(0xEF, 0x44, 0x44)
MUTED   = RGBColor(0x6B, 0x72, 0x80)
BLUE_C  = RGBColor(0x4A, 0x9E, 0xED)
AMBER   = RGBColor(0xF5, 0x9E, 0x0B)
PURPLE  = RGBColor(0x8B, 0x5C, 0xF6)
NAVY    = RGBColor(0x1E, 0x3A, 0x5F)

MO = "#FF5C35"; MD = "#141423"; MG = "#00C853"
MR = "#EF4444"; MB = "#4A9EED"; MA = "#F59E0B"
MP = "#8B5CF6"; MM = "#6B7280"; MLG = "#F5F5F7"
MN = "#1E3A5F"; MC = "#06B6D4"

W, H = Inches(13.33), Inches(7.5)

# ── Data ──────────────────────────────────────────────────────────────────────
def load():
    path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                        "Deel Marketing Analytics Challenge - raw_data.csv")
    df = pd.read_csv(path, sep=";", encoding="utf-8")
    df.columns = df.columns.str.strip().str.replace("\xa0", " ")
    df["Reporting Date"] = pd.to_datetime(df["Reporting Date"], dayfirst=True)

    # Cost: strip $ and thousands commas (US format: $160,756 → 160756)
    df["Cost"] = pd.to_numeric(df["$ Cost"].str.replace(r"[$,]", "", regex=True), errors="coerce").fillna(0)

    for c in ["# Prospects","# Marketing Qualified Prospects","# Demo Prospects",
              "# Sales Qualified Opportunities","# Opportunities Won",
              "# Impressions","# Clicks"]:
        df[c] = pd.to_numeric(df[c], errors="coerce").fillna(0)
    df["YM"]       = df["Reporting Date"].dt.to_period("M").astype(str)
    df["MonthNum"] = df["Reporting Date"].dt.month
    df["Year"]     = df["Reporting Date"].dt.year
    df["Cycle"]    = df["Reporting Date"].apply(
        lambda d: "Y1 (Jul23-Jun24)"
        if (d.year==2023 and d.month>=7) or (d.year==2024 and d.month<=6)
        else ("Y2 (Jul24-Jun25)"
              if (d.year==2024 and d.month>=7) or (d.year==2025 and d.month<=6)
              else "other"))
    return df

# ── Matplotlib theme ──────────────────────────────────────────────────────────
def mpl_theme(ax, title="", xlabel="", ylabel="", ytick_fmt=None):
    ax.set_facecolor(MLG)
    ax.figure.patch.set_facecolor("white")
    ax.spines[["top","right"]].set_visible(False)
    ax.spines[["left","bottom"]].set_color("#D1D5DB")
    ax.tick_params(colors="#374151", labelsize=8)
    if title:
        ax.set_title(title, fontsize=10, fontweight="bold", color=MD, pad=8)
    if xlabel: ax.set_xlabel(xlabel, fontsize=8, color=MM)
    if ylabel: ax.set_ylabel(ylabel, fontsize=8, color=MM)
    ax.grid(axis="y", color="#E5E7EB", linewidth=0.6, linestyle="--")
    if ytick_fmt == "M":
        ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x,_: f"${x/1e6:.0f}M"))
    elif ytick_fmt == "K":
        ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x,_: f"${x/1e3:.0f}K"))

def fig2img(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                facecolor="white", edgecolor="none")
    buf.seek(0)
    plt.close(fig)
    return buf

# ── PPT helpers ───────────────────────────────────────────────────────────────
def add_slide(prs, layout=6):
    return prs.slides.add_slide(prs.slide_layouts[layout])

def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def txt(slide, text, x, y, w, h, size=14, bold=False, color=DARK,
        align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb

def rect(slide, x, y, w, h, fill, line=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    return shape

def img(slide, buf, x, y, w, h=None):
    if h:
        slide.shapes.add_picture(buf, x, y, w, h)
    else:
        slide.shapes.add_picture(buf, x, y, w)

def orange_bar(slide, h=Inches(0.07)):
    rect(slide, 0, H - h, W, h, ORANGE)

def header_bar(slide, title, subtitle=""):
    rect(slide, 0, 0, W, Inches(1.05), DARK)
    txt(slide, title, Inches(0.4), Inches(0.10), Inches(12.5), Inches(0.6),
        size=26, bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle, Inches(0.4), Inches(0.68), Inches(12.5), Inches(0.32),
            size=12, color=ORANGE)

def kpi_card(slide, x, y, w, h, label, value, delta="", accent=ORANGE):
    rect(slide, x, y, w, h, LGRAY)
    rect(slide, x, y, w, Inches(0.06), accent)
    txt(slide, label, x+Inches(0.15), y+Inches(0.12), w-Inches(0.3), Inches(0.28),
        size=9, color=MUTED)
    txt(slide, value, x+Inches(0.15), y+Inches(0.42), w-Inches(0.3), Inches(0.5),
        size=20, bold=True, color=DARK)
    if delta:
        txt(slide, delta, x+Inches(0.15), y+Inches(0.95), w-Inches(0.3), Inches(0.28),
            size=9, bold=True, color=GREEN)

def insight_box(slide, x, y, w, h, text, accent=ORANGE):
    rect(slide, x, y, w, h, RGBColor(0xFF,0xF0,0xEC))
    rect(slide, x, y, Inches(0.07), h, accent)
    txt(slide, text, x+Inches(0.18), y+Inches(0.08), w-Inches(0.28), h-Inches(0.1),
        size=10.5, color=DARK)

# ── CHART FUNCTIONS ───────────────────────────────────────────────────────────
def chart_funnel_waterfall(df):
    stages = ["Prospects","MQPs","Demos","SQOs","Won"]
    vals = [df["# Prospects"].sum(),
            df["# Marketing Qualified Prospects"].sum(),
            df["# Demo Prospects"].sum(),
            df["# Sales Qualified Opportunities"].sum(),
            df["# Opportunities Won"].sum()]
    fig, ax = plt.subplots(figsize=(6, 3.6))
    bar_h = 0.55
    y_pos = np.arange(len(stages))[::-1]
    max_v = vals[0]
    colors = [MO, MB, MB, MA, MG]
    for i, (stage, v, y, c) in enumerate(zip(stages, vals, y_pos, colors)):
        ax.barh(y, v/max_v, bar_h, color=c, alpha=0.85)
        ax.text(v/max_v+0.01, y, f"{v:,.0f}", va="center", fontsize=8.5,
                color=MD, fontweight="bold")
        if i > 0:
            drop_pct = (1 - vals[i]/vals[i-1])*100
            ax.text(-0.05, y+bar_h/2+0.05, f"↓{drop_pct:.0f}%", va="bottom",
                    ha="right", fontsize=7, color=MR)
    ax.set_yticks(y_pos)
    ax.set_yticklabels(stages, fontsize=9)
    ax.set_xlim(-0.12, 1.28)
    ax.set_xticks([])
    mpl_theme(ax, "Global Sales Funnel — Full Period")
    ax.spines["left"].set_visible(False)
    fig.tight_layout()
    return fig2img(fig)

def chart_funnel_over_time(df):
    m = (df.groupby("YM").agg(
        demo=("# Demo Prospects","sum"),
        sqo=("# Sales Qualified Opportunities","sum"),
        won=("# Opportunities Won","sum"),
        cost=("Cost","sum")).sort_values("YM").reset_index())
    x = range(len(m))
    labels = [v[2:7] for v in m["YM"]]
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(11, 3.4))
    ax1.plot(x, m["demo"], color=MB, lw=2, marker="o", ms=3, label="Demos")
    ax1.plot(x, m["sqo"],  color=MA, lw=2, marker="s", ms=3, label="SQOs")
    ax1.plot(x, m["won"],  color=MG, lw=2.5, marker="^", ms=4, label="Won")
    ax1.set_xticks(list(x)[::3])
    ax1.set_xticklabels(labels[::3], rotation=45, ha="right", fontsize=7)
    mpl_theme(ax1, "Bottom-Funnel Volumes Over Time")
    ax1.legend(fontsize=8)
    peak_i = m["won"].idxmax()
    ax1.annotate(f"Peak: {m['won'].iloc[peak_i]:.0f}",
                 xy=(peak_i, m["won"].iloc[peak_i]),
                 xytext=(peak_i-3, m["won"].iloc[peak_i]+40),
                 fontsize=7, color=MG,
                 arrowprops=dict(arrowstyle="->", color=MG, lw=0.8))

    ax2.bar(x, m["cost"], color=MO, alpha=0.6, label="Spend ($M)")
    ax2b = ax2.twinx()
    ax2b.plot(x, m["won"], color=MN, lw=2, marker="o", ms=3, label="Won")
    ax2b.spines[["top","right"]].set_visible(False)
    ax2b.tick_params(colors="#374151", labelsize=7)
    ax2.set_xticks(list(x)[::3])
    ax2.set_xticklabels(labels[::3], rotation=45, ha="right", fontsize=7)
    mpl_theme(ax2, "Monthly Spend vs Won", ytick_fmt="M")
    h1,l1 = ax2.get_legend_handles_labels()
    h2,l2 = ax2b.get_legend_handles_labels()
    ax2.legend(h1+h2, l1+l2, fontsize=7, loc="upper left")
    fig.tight_layout(pad=1.5)
    return fig2img(fig)

def chart_impressions_clicks(df):
    m = (df.groupby("YM").agg(
        imp=("# Impressions","sum"),
        clk=("# Clicks","sum")).sort_values("YM").reset_index())
    x = range(len(m))
    labels = [v[2:7] for v in m["YM"]]
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(11, 3.2))
    ax1.fill_between(x, m["imp"]/1e6, color=MO, alpha=0.25)
    ax1.plot(x, m["imp"]/1e6, color=MO, lw=2, label="Impressions (M)")
    ax1.set_xticks(list(x)[::3])
    ax1.set_xticklabels(labels[::3], rotation=45, ha="right", fontsize=7)
    ax1.yaxis.set_major_formatter(mticker.FuncFormatter(lambda v,_: f"{v:.0f}M"))
    mpl_theme(ax1, "Impressions Over Time")
    peak_imp = m["imp"].idxmax()
    ax1.annotate(f"Jan'24 spike\n{m['imp'].iloc[peak_imp]/1e6:.0f}M",
                 xy=(peak_imp, m["imp"].iloc[peak_imp]/1e6),
                 xytext=(peak_imp+1, m["imp"].iloc[peak_imp]/1e6*0.82),
                 fontsize=7, color=MO,
                 arrowprops=dict(arrowstyle="->", color=MO, lw=0.8))

    ax2.fill_between(x, m["clk"]/1e6, color=MB, alpha=0.25)
    ax2.plot(x, m["clk"]/1e6, color=MB, lw=2, label="Clicks (M)")
    ax2.set_xticks(list(x)[::3])
    ax2.set_xticklabels(labels[::3], rotation=45, ha="right", fontsize=7)
    ax2.yaxis.set_major_formatter(mticker.FuncFormatter(lambda v,_: f"{v:.1f}M"))
    mpl_theme(ax2, "Clicks Over Time")
    fig.tight_layout()
    return fig2img(fig)

def chart_seasonality(df):
    fiscal_order = list(range(7,13)) + list(range(1,7))
    month_lbl = {m: ["Jan","Feb","Mar","Apr","May","Jun",
                      "Jul","Aug","Sep","Oct","Nov","Dec"][m-1] for m in range(1,13)}
    cyc = df[df["Cycle"] != "other"].copy()
    yoy = (cyc.groupby(["Cycle","MonthNum"])
           .agg(won=("# Opportunities Won","sum"), cost=("Cost","sum"))
           .reset_index())
    yoy["fp"] = yoy["MonthNum"].map({m:i+1 for i,m in enumerate(fiscal_order)})
    yoy = yoy.sort_values(["Cycle","fp"])
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(11, 3.2))
    for cyc_name, color, ls in [("Y1 (Jul23-Jun24)", MB, "--"), ("Y2 (Jul24-Jun25)", MO, "-")]:
        d = yoy[yoy["Cycle"]==cyc_name]
        ax1.plot(d["fp"], d["won"], color=color, lw=2.2, linestyle=ls,
                 marker="o", ms=5, label=cyc_name)
    ax1.set_xticks(range(1,13))
    ax1.set_xticklabels([month_lbl[m] for m in fiscal_order], fontsize=8)
    mpl_theme(ax1, "Won by Fiscal Month — YoY")
    ax1.legend(fontsize=8)
    y2 = yoy[yoy["Cycle"]=="Y2 (Jul24-Jun25)"]
    pk = y2.loc[y2["won"].idxmax()]
    ax1.annotate(f"Peak: {pk['won']:.0f}", xy=(pk["fp"], pk["won"]),
                 xytext=(pk["fp"]-1.5, pk["won"]+25),
                 fontsize=7.5, color=MG,
                 arrowprops=dict(arrowstyle="->", color=MG, lw=0.8))
    for cyc_name, color, ls in [("Y1 (Jul23-Jun24)", MB, "--"), ("Y2 (Jul24-Jun25)", MO, "-")]:
        d = yoy[yoy["Cycle"]==cyc_name]
        ax2.plot(d["fp"], d["cost"]/1e6, color=color, lw=2.2, linestyle=ls,
                 marker="o", ms=5, label=cyc_name)
    ax2.set_xticks(range(1,13))
    ax2.set_xticklabels([month_lbl[m] for m in fiscal_order], fontsize=8)
    ax2.yaxis.set_major_formatter(mticker.FuncFormatter(lambda v,_: f"${v:.1f}M"))
    mpl_theme(ax2, "Investment by Fiscal Month — YoY")
    ax2.legend(fontsize=8)
    fig.tight_layout()
    return fig2img(fig)

def chart_team_funnel(df):
    by_t = (df.groupby("Marketing Team")
            .agg(prospects=("# Prospects","sum"),
                 mqp=("# Marketing Qualified Prospects","sum"),
                 demo=("# Demo Prospects","sum"),
                 sqo=("# Sales Qualified Opportunities","sum"),
                 won=("# Opportunities Won","sum"),
                 cost=("Cost","sum")).reset_index())
    safe = lambda a,b: np.where(b>0, a/b*100, 0)
    by_t["E2E%"] = safe(by_t["won"], by_t["prospects"]).round(2)
    by_t["CPA"]  = np.where(by_t["cost"]>0, by_t["cost"]/by_t["won"].replace(0,np.nan), 0)
    by_t = by_t.sort_values("E2E%", ascending=True)

    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(11, 4.8))
    colors1 = [MG if v >= 2 else (MA if v >= 0.5 else MR) for v in by_t["E2E%"]]
    bars = ax1.barh(by_t["Marketing Team"], by_t["E2E%"], color=colors1, alpha=0.85, height=0.55)
    for bar, v in zip(bars, by_t["E2E%"]):
        ax1.text(bar.get_width()+0.02, bar.get_y()+bar.get_height()/2,
                 f"{v:.2f}%", va="center", fontsize=8.5, color=MD, fontweight="bold")
    mpl_theme(ax1, "End-to-End Rate (Won / Prospects) by Team")
    ax1.set_xlabel("E2E Rate (%)", fontsize=8, color=MM)

    cb = by_t[by_t["CPA"]>0].sort_values("CPA")
    med_cpa = cb["CPA"].median()
    colors2 = [MG if v<=med_cpa else (MA if v<=med_cpa*3 else MR) for v in cb["CPA"]]
    bars2 = ax2.barh(cb["Marketing Team"], cb["CPA"], color=colors2, alpha=0.85, height=0.55)
    for bar, v in zip(bars2, cb["CPA"]):
        ax2.text(bar.get_width()+0.3, bar.get_y()+bar.get_height()/2,
                 f"${v:,.0f}", va="center", fontsize=8.5, color=MD, fontweight="bold")
    ax2.xaxis.set_major_formatter(mticker.FuncFormatter(lambda v,_: f"${v:,.0f}"))
    mpl_theme(ax2, "CPA by Team (Cost per Won Opp)")
    ax2.set_xlabel("CPA ($)", fontsize=8, color=MM)
    fig.tight_layout(pad=2.0)
    return fig2img(fig)

def chart_channel_conversion(df):
    channels = ["paid_search", "event", "regional"]
    ch = (df[df["Marketing Program"].isin(channels)]
          .groupby("Marketing Program")
          .agg(imp=("# Impressions","sum"), clk=("# Clicks","sum"),
               demo=("# Demo Prospects","sum"), sqo=("# Sales Qualified Opportunities","sum"),
               won=("# Opportunities Won","sum"), cost=("Cost","sum")).reset_index())
    ch["CTR%"]      = np.where(ch["imp"]>0, ch["clk"]/ch["imp"]*100, np.nan)
    ch["Clk→Demo%"] = np.where(ch["clk"]>0, ch["demo"]/ch["clk"]*100, np.nan)
    ch["Demo→SQO%"] = np.where(ch["demo"]>0, ch["sqo"]/ch["demo"]*100, np.nan)
    ch["SQO→Won%"]  = np.where(ch["sqo"]>0, ch["won"]/ch["sqo"]*100, np.nan)
    ch["CPA"]       = np.where(ch["won"]>0, ch["cost"]/ch["won"], np.nan)

    metrics = ["CTR%","Clk→Demo%","Demo→SQO%","SQO→Won%"]
    colors  = [MB, MP, MA, MG]
    x = np.arange(len(channels)); bw = 0.2

    fig, ax = plt.subplots(figsize=(9, 3.5))
    for i, (m, c) in enumerate(zip(metrics, colors)):
        vals = ch[m].fillna(0).tolist()
        bars = ax.bar(x + i*bw, vals, bw, label=m, color=c, alpha=0.85)
        for bar, v in zip(bars, vals):
            if v > 0:
                ax.text(bar.get_x()+bar.get_width()/2, bar.get_height()+0.3,
                        f"{v:.2f}%" if m=="CTR%" else f"{v:.1f}%",
                        ha="center", va="bottom", fontsize=6.5, color=MD)
    ax.set_xticks(x + bw*1.5)
    ax.set_xticklabels(channels, fontsize=9)
    mpl_theme(ax, "Conversion Rates: Paid Search vs Event vs Regional")
    ax.legend(fontsize=7.5, ncol=2)
    ax.set_ylabel("Rate (%)", fontsize=8, color=MM)
    fig.tight_layout()
    return fig2img(fig)

def chart_paid_search_brand(df):
    ps = (df[df["Marketing Program"]=="paid_search"].groupby("YM")
          .agg(cost=("Cost","sum"), won=("# Opportunities Won","sum"),
               demo=("# Demo Prospects","sum"), sqo=("# Sales Qualified Opportunities","sum"))
          .reset_index().sort_values("YM"))
    bs = (df[df["Marketing Program"]=="paid_brand_search"].groupby("YM")
          .agg(cost=("Cost","sum"), won=("# Opportunities Won","sum"),
               demo=("# Demo Prospects","sum"), sqo=("# Sales Qualified Opportunities","sum"))
          .reset_index().sort_values("YM"))
    ps["CPA"] = ps["cost"]/ps["won"].replace(0, np.nan)
    bs["CPA"] = bs["cost"]/bs["won"].replace(0, np.nan)

    fig, axes = plt.subplots(2, 2, figsize=(11, 5.5))
    for ax, d, color, name in [
        (axes[0,0], ps, MB, "Paid Search — Spend vs Won"),
        (axes[0,1], bs, MG, "Brand Search — Spend vs Won"),
    ]:
        x = range(len(d))
        labels = [v[2:7] for v in d["YM"]]
        ax.bar(x, d["cost"], color=color, alpha=0.4, label="Spend ($M)")
        ax2 = ax.twinx()
        ax2.plot(x, d["won"], color=MO, lw=2, marker="o", ms=3, label="Won")
        ax2.spines[["top","right"]].set_visible(False)
        ax2.tick_params(colors="#374151", labelsize=7)
        ax.set_xticks(list(x)[::4])
        ax.set_xticklabels(labels[::4], rotation=45, ha="right", fontsize=7)
        mpl_theme(ax, name, ytick_fmt="M")
        h1,l1 = ax.get_legend_handles_labels()
        h2,l2 = ax2.get_legend_handles_labels()
        ax.legend(h1+h2, l1+l2, fontsize=6.5, loc="upper left")

    for ax, d, color, name, bmark in [
        (axes[1,0], ps, MB, "Paid Search — CPA Trend",   50000),
        (axes[1,1], bs, MG, "Brand Search — CPA Trend",  20000),
    ]:
        x = range(len(d))
        labels = [v[2:7] for v in d["YM"]]
        ax.plot(x, d["CPA"], color=color, lw=2, marker="o", ms=3)
        ax.axhline(bmark, color=MR, lw=1.2, linestyle="--", label=f"Benchmark ${bmark/1e3:.0f}K")
        ax.set_xticks(list(x)[::4])
        ax.set_xticklabels(labels[::4], rotation=45, ha="right", fontsize=7)
        ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda v,_: f"${v:,.0f}"))
        mpl_theme(ax, name)
        ax.legend(fontsize=7)
    fig.tight_layout(pad=1.5)
    return fig2img(fig)

def chart_brand_channels(df):
    br = (df[df["Marketing Team"].isin(["content","website","lifecycle"])]
          .groupby(["YM","Marketing Team"])
          .agg(won=("# Opportunities Won","sum"), pros=("# Prospects","sum"))
          .reset_index().sort_values("YM"))
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(13, 4.2))
    colors_map = {"content": MO, "website": MB, "lifecycle": MP}
    for team in ["content", "website", "lifecycle"]:
        d = br[br["Marketing Team"]==team]
        x = range(len(d))
        ax1.plot(x, d["won"], color=colors_map[team], lw=2, marker="o", ms=3, label=team.capitalize())
        ax1.set_xticks(list(x)[::4])
        ax1.set_xticklabels([v[2:7] for v in d["YM"]][::4], rotation=45, ha="right", fontsize=7)
    mpl_theme(ax1, "Zero-Cost Channels — Monthly Won")
    ax1.legend(fontsize=8)

    dc = br[br["Marketing Team"]=="content"]
    xc = range(len(dc))
    ax2.fill_between(xc, dc["pros"], color=MO, alpha=0.3)
    ax2.plot(xc, dc["pros"], color=MO, lw=2, label="Content Prospects")
    dl = br[br["Marketing Team"]=="lifecycle"]
    xl = range(len(dl))
    ax2b = ax2.twinx()
    ax2b.plot(xl, dl["won"], color=MP, lw=2, linestyle="--", label="Lifecycle Won")
    ax2b.spines[["top","right"]].set_visible(False)
    ax2b.tick_params(colors="#374151", labelsize=7)
    ax2.set_xticks(list(xc)[::4])
    ax2.set_xticklabels([v[2:7] for v in dc["YM"]][::4], rotation=45, ha="right", fontsize=7)
    mpl_theme(ax2, "Content Prospects vs Lifecycle Won")
    h1,l1 = ax2.get_legend_handles_labels()
    h2,l2 = ax2b.get_legend_handles_labels()
    ax2.legend(h1+h2, l1+l2, fontsize=7, loc="upper left")
    fig.tight_layout()
    return fig2img(fig)

def chart_reallocation(df):
    paid = (df[df["Marketing Team"]=="paid"].groupby("Marketing Program")
            .agg(cost=("Cost","sum"), won=("# Opportunities Won","sum")).reset_index())
    paid["CPA"] = np.where(paid["won"]>0, paid["cost"]/paid["won"], np.nan)
    paid = paid.dropna(subset=["CPA"]).sort_values("CPA")

    benchmark = paid["CPA"].median()  # data-driven: splits efficient vs inefficient programs
    bmark_lbl = f"${benchmark/1e3:.0f}K"
    fig, ax = plt.subplots(figsize=(9, 4.2))
    colors = [MG if v<=benchmark else MR for v in paid["CPA"]]
    bars = ax.barh(paid["Marketing Program"], paid["CPA"], color=colors, alpha=0.85, height=0.55)
    ax.axvline(benchmark, color="#9ca3af", lw=1.5, linestyle="--")
    ax.text(benchmark*1.02, -0.5, f"{bmark_lbl}\nbenchmark", color="#6b7280", fontsize=7.5, va="top")
    for bar, v in zip(bars, paid["CPA"]):
        ax.text(bar.get_width()+benchmark*0.02, bar.get_y()+bar.get_height()/2,
                f"${v:,.0f}", va="center", fontsize=8, color=MD, fontweight="bold")
    mpl_theme(ax, f"CPA by Paid Program  (green = below {bmark_lbl} median benchmark)")
    ax.set_xlabel("CPA ($)", fontsize=8, color=MM)
    ax.xaxis.set_major_formatter(mticker.FuncFormatter(lambda v,_: f"${v/1e3:.0f}K"))
    patches = [mpatches.Patch(color=MG, label=f"Scale Up (CPA ≤ {bmark_lbl})"),
               mpatches.Patch(color=MR, label=f"Review / Cut (CPA > {bmark_lbl})")]
    ax.legend(handles=patches, fontsize=7.5, loc="lower right")
    fig.tight_layout(pad=2.0)
    return fig2img(fig)

def chart_budget_impact(df):
    paid = (df[df["Marketing Team"]=="paid"].groupby("Marketing Program")
            .agg(cost=("Cost","sum"), won=("# Opportunities Won","sum")).reset_index())
    paid["CPA"] = np.where(paid["won"]>0, paid["cost"]/paid["won"], np.nan)
    paid = paid.dropna(subset=["CPA"]).sort_values("CPA").reset_index(drop=True)

    progs = paid["Marketing Program"].tolist()
    cpas  = paid["CPA"].tolist()
    costs = paid["cost"].tolist()
    n = len(progs)

    # Bottom third by CPA = scale up; top third = cut
    n_low  = max(n // 3, 1)
    n_high = max(n // 3, 1)
    low_set  = set(progs[:n_low])
    high_set = set(progs[n-n_high:])

    cut_amount = sum(c for p,c in zip(progs,costs) if p in high_set) * 0.5
    add_per    = cut_amount / max(len(low_set), 1)

    opt_costs = {}
    for p, c in zip(progs, costs):
        if p in high_set:
            opt_costs[p] = c * 0.5
        elif p in low_set:
            opt_costs[p] = c + add_per
        else:
            opt_costs[p] = c

    cpa_dict = dict(zip(progs, cpas))
    won_curr = sum(c/cpa_dict[p] for p,c in zip(progs,costs) if cpa_dict[p]>0)
    won_opt  = sum(opt_costs[p]/cpa_dict[p] for p in progs if cpa_dict[p]>0)

    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(10, 3.8))
    width = 0.35
    x = np.arange(n)
    bar_colors_curr = [MG if p in low_set else (MR if p in high_set else MB) for p in progs]
    ax1.bar(x - width/2, costs,              width, label="Current ($)",   color=MB,  alpha=0.65)
    ax1.bar(x + width/2, [opt_costs[p] for p in progs], width, label="Optimized ($)", color=MO, alpha=0.65)
    ax1.set_xticks(x)
    ax1.set_xticklabels([p.replace("paid_","") for p in progs], rotation=45, ha="right", fontsize=7)
    ax1.yaxis.set_major_formatter(mticker.FuncFormatter(lambda v,_: f"${v/1e6:.0f}M"))
    mpl_theme(ax1, "Budget Allocation: Current vs Optimized ($)")
    ax1.legend(fontsize=8)

    won_vals   = [won_curr, won_opt]
    bar_colors = [MB, MG]
    bars = ax2.bar(["Current","Optimized"], won_vals, color=bar_colors, alpha=0.85, width=0.45)
    for bar, v in zip(bars, won_vals):
        ax2.text(bar.get_x()+bar.get_width()/2, bar.get_height()+1,
                 f"{v:,.0f}", ha="center", va="bottom", fontsize=11, fontweight="bold", color=MD)
    if won_curr > 0:
        delta = won_opt - won_curr
        sign  = "+" if delta >= 0 else ""
        ax2.text(1, won_opt*0.45, f"{sign}{delta:,.0f}\n({sign}{delta/won_curr*100:.0f}%)",
                 ha="center", va="center", fontsize=12, fontweight="bold", color=MG)
    mpl_theme(ax2, "Projected Won: Current vs Optimized")
    ax2.set_ylabel("Won Opportunities", fontsize=8, color=MM)
    fig.tight_layout(pad=1.5)
    return fig2img(fig)

def chart_cpm_ctr(df):
    paid = df[df["Marketing Team"]=="paid"]
    imp = (paid[paid["# Impressions"]>0]
           .groupby("Marketing Program")
           .agg(impressions=("# Impressions","sum"), clicks=("# Clicks","sum"),
                cost=("Cost","sum"), won=("# Opportunities Won","sum")).reset_index())
    imp["CTR"] = imp["clicks"]/imp["impressions"]*100
    imp["CPM"] = imp["cost"]/imp["impressions"]*1000
    imp["CPA"] = np.where(imp["won"]>0, imp["cost"]/imp["won"], np.nan)
    imp = imp.dropna(subset=["CPA"])

    med_cpa = imp["CPA"].median()
    fig, ax = plt.subplots(figsize=(9, 4.2))
    scatter_colors = [MG if v<=med_cpa else (MA if v<=med_cpa*3 else MR) for v in imp["CPA"]]
    # log-scale bubble sizes so paid_awareness ($2.1M CPA) doesn't dwarf the rest
    log_cpa = np.log1p(imp["CPA"])
    sizes = (log_cpa / log_cpa.max() * 1000 + 100).tolist()
    ax.scatter(imp["CPM"], imp["CTR"], s=sizes, c=scatter_colors, alpha=0.65, zorder=5,
               edgecolors="white", linewidths=0.5)
    for _, row in imp.iterrows():
        name = row["Marketing Program"].replace("paid_","").replace("_leadbuy_listing","")
        cpm, ctr = row["CPM"], row["CTR"]
        # custom offsets to avoid overlap in the clustered bottom-left region
        if cpm > 1000:                        # brand_search — far right, pull label left
            offset = (-75, 6)
        elif cpm < 20 and ctr > 0.3:         # social — push up
            offset = (6, 10)
        elif cpm < 20 and ctr <= 0.3:        # listicle — push down
            offset = (6, -16)
        elif cpm < 50 and ctr < 0.2:         # awareness — pull left/up
            offset = (-70, 10)
        else:
            offset = (6, 6)
        ax.annotate(name, (row["CPM"], row["CTR"]),
                    xytext=offset, textcoords="offset points", fontsize=8, color=MD,
                    fontweight="bold")
    med_ctr = imp["CTR"].median(); med_cpm = imp["CPM"].median()
    ax.axhline(med_ctr, color="#9ca3af", lw=1.2, linestyle="--",
               label=f"Median CTR {med_ctr:.2f}%")
    ax.axvline(med_cpm, color="#9ca3af", lw=1.2, linestyle=":",
               label=f"Median CPM ${med_cpm:.2f}")
    ax.set_xlabel("CPM ($) — lower is better", fontsize=9, color=MM)
    ax.set_ylabel("CTR (%) — higher is better", fontsize=9, color=MM)
    mpl_theme(ax, "CPM vs CTR Matrix  (bubble size = CPA)")
    ax.set_xlim(left=-100, right=2000)
    ax.set_ylim(bottom=-2, top=8)
    ax.xaxis.set_major_formatter(mticker.FuncFormatter(lambda v,_: f"${v:,.0f}"))
    patches = [mpatches.Patch(color=MG, label=f"Low CPA (≤${med_cpa/1e3:.0f}K)"),
               mpatches.Patch(color=MA, label=f"Mid CPA (≤${med_cpa*3/1e3:.0f}K)"),
               mpatches.Patch(color=MR, label=f"High CPA (>${med_cpa*3/1e3:.0f}K)")]
    ax.legend(handles=patches, fontsize=8, loc="upper right")
    fig.tight_layout()
    return fig2img(fig)

# ── SLIDE BUILDERS ────────────────────────────────────────────────────────────
def slide_cover(prs):
    s = add_slide(prs)
    bg(s, DARK)
    rect(s, 0, 0, Inches(0.18), H, ORANGE)
    rect(s, 0, H-Inches(0.06), W, Inches(0.06), ORANGE)
    txt(s, "Deel", Inches(0.5), Inches(0.9), Inches(4), Inches(0.7),
        size=44, bold=True, color=ORANGE)
    txt(s, "Marketing Analytics Challenge",
        Inches(0.5), Inches(1.65), Inches(10), Inches(0.7),
        size=30, bold=True, color=WHITE)
    txt(s, "Executive Presentation  |  Tasks 1, 2 & 3",
        Inches(0.5), Inches(2.45), Inches(10), Inches(0.45),
        size=15, color=RGBColor(0xCC,0xCC,0xCC))
    txt(s, "Period: Jul 2023 – Jun 2025   |   11,952 records   |   7 teams   |   17 programs",
        Inches(0.5), Inches(3.1), Inches(10), Inches(0.4),
        size=12, color=MUTED)
    txt(s, "Tools: Python · Pandas · NumPy · Plotly · Streamlit · Matplotlib · python-pptx",
        Inches(0.5), Inches(3.55), Inches(10), Inches(0.4),
        size=11, color=MUTED)
    for i, (title, sub) in enumerate([
        ("Task 1", "Sales Funnel Analysis"),
        ("Task 2", "Paid & Brand Channels"),
        ("Task 3", "Scenario Analysis"),
    ]):
        x = Inches(0.5 + i * 4.1)
        rect(s, x, Inches(4.5), Inches(3.8), Inches(1.55), RGBColor(0x22,0x22,0x38))
        rect(s, x, Inches(4.5), Inches(3.8), Inches(0.06), ORANGE)
        txt(s, title, x+Inches(0.2), Inches(4.65), Inches(3.4), Inches(0.45),
            size=20, bold=True, color=ORANGE)
        txt(s, sub, x+Inches(0.2), Inches(5.2), Inches(3.4), Inches(0.45),
            size=13, color=WHITE)
    txt(s, "Deel Marketing Intelligence  |  2025",
        Inches(0.5), Inches(6.7), Inches(8), Inches(0.4), size=10, color=MUTED)

def slide_audience(prs):
    s = add_slide(prs)
    bg(s, WHITE)
    header_bar(s, "Who You're Presenting To", "Two senior stakeholders with distinct focus areas")
    orange_bar(s)
    for i, (role, icon_label, focus, q1, q2, q3, c) in enumerate([
        ("Paid Marketing Director", "ROI FOCUS",
         "Cares about short-term ROI and performance of paid spend.",
         "Which channels generate the best CPA?",
         "Where is budget being wasted?",
         "How do we grow Won at lower cost?", ORANGE),
        ("Head of Marketing", "BRAND FOCUS",
         "Wants to prove long-term impact of awareness and brand-building.",
         "Are organic channels growing over time?",
         "What is the implied value of zero-cost channels?",
         "How do we measure brand budget impact?", BLUE_C),
    ]):
        x = Inches(0.35 + i * 6.55)
        rect(s, x, Inches(1.15), Inches(6.25), Inches(5.9), LGRAY)
        rect(s, x, Inches(1.15), Inches(6.25), Inches(0.55), c)
        txt(s, role, x+Inches(0.15), Inches(1.2), Inches(5.9), Inches(0.5),
            size=16, bold=True, color=WHITE)
        txt(s, focus, x+Inches(0.2), Inches(1.88), Inches(5.85), Inches(0.4),
            size=11, color=DARK)
        txt(s, "Key Questions:", x+Inches(0.2), Inches(2.38), Inches(5.85), Inches(0.32),
            size=10, bold=True, color=DARK)
        for j, q in enumerate([q1, q2, q3]):
            txt(s, f"→  {q}", x+Inches(0.2), Inches(2.75+j*0.55), Inches(5.85), Inches(0.5),
                size=10, color=MUTED)

def slide_global_kpis(prs, df):
    s = add_slide(prs)
    bg(s, WHITE)
    header_bar(s, "Portfolio Overview — Key Metrics at a Glance",
               "Full period Jul 2023 – Jun 2025  |  All channels combined")
    orange_bar(s)

    tc   = df["Cost"].sum()
    tw   = df["# Opportunities Won"].sum()
    tp   = df["# Prospects"].sum()
    cpa  = tc / tw if tw else 0
    e2e  = tw / tp * 100 if tp else 0
    tclk = df["# Clicks"].sum()
    timp = df["# Impressions"].sum()
    tsqo = df["# Sales Qualified Opportunities"].sum()
    tdemo= df["# Demo Prospects"].sum()
    tmqp = df["# Marketing Qualified Prospects"].sum()

    kpis1 = [
        ("Total Investment", f"${tc/1e6:.0f}M", ORANGE),
        ("Total Won Opps.", f"{tw:,.0f}", GREEN),
        ("Avg. Global CPA", f"${cpa:,.0f}", BLUE_C),
        ("E2E Conv. Rate", f"{e2e:.2f}%", PURPLE),
    ]
    kpis2 = [
        ("Total Impressions", f"{timp/1e9:.1f}B", NAVY),
        ("Total Clicks", f"{tclk/1e6:.1f}M", BLUE_C),
        ("Total Demos", f"{tdemo:,.0f}", AMBER),
        ("Total SQOs", f"{tsqo:,.0f}", ORANGE),
    ]

    w = Inches(3.1); gap = Inches(0.12)
    for i, (label, value, c) in enumerate(kpis1):
        kpi_card(s, Inches(0.3)+i*(w+gap), Inches(1.2), w, Inches(1.3), label, value, accent=c)
    for i, (label, value, c) in enumerate(kpis2):
        kpi_card(s, Inches(0.3)+i*(w+gap), Inches(2.65), w, Inches(1.3), label, value, accent=c)

    buf = chart_funnel_waterfall(df)
    img(s, buf, Inches(0.3), Inches(4.05), Inches(5.8), Inches(3.1))

    rect(s, Inches(6.3), Inches(4.05), Inches(6.8), Inches(3.1), LGRAY)
    txt(s, "Funnel Drop-off Summary", Inches(6.5), Inches(4.15), Inches(6.5), Inches(0.35),
        size=12, bold=True, color=DARK)
    funnel_data = [
        ("Prospects → MQPs",   f"{tmqp/tp*100:.1f}%",  "filter applied at top of funnel"),
        ("MQPs → Demos",       f"{tdemo/tmqp*100:.1f}%" if tmqp else "n/a", "quality gate"),
        ("Demos → SQOs",       f"{tsqo/tdemo*100:.1f}%", "sales qualification"),
        ("SQOs → Won",         f"{tw/tsqo*100:.1f}%",   "deal closure"),
        ("Overall E2E",        f"{e2e:.2f}%",           "end-to-end conversion"),
    ]
    for j, (stage, rate, note) in enumerate(funnel_data):
        y = Inches(4.55) + j * Inches(0.5)
        rect(s, Inches(6.4), y, Inches(0.06), Inches(0.35), ORANGE)
        txt(s, stage, Inches(6.55), y, Inches(2.5), Inches(0.35), size=10, bold=True, color=DARK)
        txt(s, rate, Inches(9.1), y, Inches(0.9), Inches(0.35), size=13, bold=True, color=ORANGE,
            align=PP_ALIGN.RIGHT)
        txt(s, note, Inches(10.1), y, Inches(2.9), Inches(0.35), size=9, color=MUTED)

def slide_task_divider(prs, task_num, task_title, subtitle):
    s = add_slide(prs)
    bg(s, DARK)
    rect(s, 0, 0, Inches(0.18), H, ORANGE)
    rect(s, 0, H-Inches(0.06), W, Inches(0.06), ORANGE)
    txt(s, f"TASK {task_num}", Inches(0.5), Inches(2.1), Inches(4), Inches(0.6),
        size=14, bold=True, color=ORANGE)
    txt(s, task_title, Inches(0.5), Inches(2.75), Inches(12), Inches(0.85),
        size=34, bold=True, color=WHITE)
    txt(s, subtitle, Inches(0.5), Inches(3.68), Inches(12), Inches(0.5),
        size=15, color=RGBColor(0xCC,0xCC,0xCC))

def slide_top_funnel(prs, df):
    s = add_slide(prs)
    bg(s, WHITE)
    header_bar(s, "Task 1 — Top-Funnel: Impressions & Clicks Over Time",
               "Awareness stage metrics  |  Jul 2023 – Jun 2025")
    orange_bar(s)
    buf = chart_impressions_clicks(df)
    img(s, buf, Inches(0.3), Inches(1.12), Inches(9.2), Inches(3.8))

    timp = df["# Impressions"].sum()
    tclk = df["# Clicks"].sum()
    ctr  = tclk/timp*100 if timp else 0

    facts = [
        ("Total Impressions", f"{timp/1e9:.1f}B", NAVY),
        ("Total Clicks", f"{tclk/1e6:.1f}M", BLUE_C),
        ("Overall CTR", f"{ctr:.2f}%", ORANGE),
    ]
    for i, (label, val, c) in enumerate(facts):
        kpi_card(s, Inches(9.65), Inches(1.12)+i*Inches(1.42), Inches(3.4), Inches(1.2),
                 label, val, accent=c)

    observations = [
        "Jan 2024 spike: 820M impressions vs ~380M avg — likely a large awareness campaign launch.",
        "CTR of 3.64% for paid_search is strong; paid_social at 0.47% signals poor audience fit.",
        "Click volume grew steadily from Jul 2023 to Oct 2024, then plateaued — market saturation signal.",
        "Impressions and clicks are concentrated in the paid team (~6K of 11.9K rows); organic channels have no impression data.",
    ]
    rect(s, Inches(0.3), Inches(5.1), Inches(12.7), Inches(2.0), LGRAY)
    txt(s, "Observations", Inches(0.5), Inches(5.2), Inches(12.3), Inches(0.32),
        size=11, bold=True, color=DARK)
    for j, obs in enumerate(observations):
        col = j % 2; row_i = j // 2
        x = Inches(0.5 + col * 6.4); y = Inches(5.55) + row_i * Inches(0.45)
        rect(s, x-Inches(0.12), y+Inches(0.06), Inches(0.06), Inches(0.22), ORANGE)
        txt(s, obs, x, y, Inches(6.1), Inches(0.4), size=9.5, color=DARK)

def slide_bottom_funnel(prs, df):
    s = add_slide(prs)
    bg(s, WHITE)
    header_bar(s, "Task 1 — Bottom-Funnel: Demos, SQOs & Won Over Time",
               "Conversion outcomes  |  Jul 2023 – Jun 2025")
    orange_bar(s)
    buf = chart_funnel_over_time(df)
    img(s, buf, Inches(0.3), Inches(1.12), Inches(12.7), Inches(4.0))

    insights = [
        ("Q4 Surge (Oct–Dec)", "Enterprise year-end budget decisions accelerate; Y2 peaked at 622 Won in May. Nov-Dec show consistent uplift across both cycles."),
        ("Jan–Feb Dip", "Budget resets slow deal closures. Jan 2025 dropped to ~350 Won despite high spend — CPA spikes to $50–65K in this period."),
        ("Upward Trend", "Total Won grew from ~285/mo (Y1 avg) to ~580/mo (Y2 avg) — a 2x improvement driven by paid scale-up and lifecycle growth."),
    ]
    for i, (title, body) in enumerate(insights):
        x = Inches(0.3 + i * 4.35)
        rect(s, x, Inches(5.25), Inches(4.15), Inches(1.85), LGRAY)
        rect(s, x, Inches(5.25), Inches(4.15), Inches(0.06), ORANGE)
        txt(s, title, x+Inches(0.15), Inches(5.35), Inches(3.95), Inches(0.35),
            size=11, bold=True, color=DARK)
        txt(s, body, x+Inches(0.15), Inches(5.75), Inches(3.95), Inches(1.2),
            size=9.5, color=MUTED)

def slide_seasonality(prs, df):
    s = add_slide(prs)
    bg(s, WHITE)
    header_bar(s, "Task 1 — Seasonality Patterns by Fiscal Cycle",
               "Fiscal year: Jul (Month 1) → Jun (Month 12)")
    orange_bar(s)
    buf = chart_seasonality(df)
    img(s, buf, Inches(0.3), Inches(1.12), Inches(8.6), Inches(3.8))

    rect(s, Inches(9.1), Inches(1.12), Inches(4.0), Inches(5.95), LGRAY)
    txt(s, "Seasonality Drivers", Inches(9.25), Inches(1.22), Inches(3.8), Inches(0.35),
        size=12, bold=True, color=DARK)
    obs = [
        ("Oct–Dec Peak", "Enterprise fiscal year-end. Companies accelerate buying decisions before Dec close."),
        ("May–Jun Peak (Y2)", "Strong Q4 performance repeated — suggests maturing sales cycle management."),
        ("Jan–Mar Slow", "Budget resets. Decision-makers pause. CPA rises 30–50% vs Q4 levels."),
        ("Jan'24 Spend Spike", "Investment jumped 3x in Jan 2024 — aggressive push against Jan dip; limited Won impact."),
        ("Y2 Consistently Higher", "Won per month up 2x vs Y1 across all months, driven by lifecycle and paid optimization."),
    ]
    y0 = Inches(1.65)
    for (title, body) in obs:
        rect(s, Inches(9.2), y0, Inches(0.05), Inches(0.3), ORANGE)
        txt(s, title, Inches(9.35), y0, Inches(3.6), Inches(0.28),
            size=10, bold=True, color=DARK)
        txt(s, body, Inches(9.35), y0+Inches(0.28), Inches(3.6), Inches(0.48),
            size=9, color=MUTED)
        y0 += Inches(0.88)

def slide_team_funnel(prs, df):
    s = add_slide(prs)
    bg(s, WHITE)
    header_bar(s, "Task 1 — Funnel Efficiency by Marketing Team",
               "Which teams convert best? E2E Rate = Won / Prospects")
    orange_bar(s)
    buf = chart_team_funnel(df)
    img(s, buf, Inches(0.3), Inches(1.12), Inches(9.0), Inches(4.5))

    by_t = (df.groupby("Marketing Team")
            .agg(pros=("# Prospects","sum"), won=("# Opportunities Won","sum"),
                 cost=("Cost","sum")).reset_index())
    by_t["E2E%"] = np.where(by_t["pros"]>0, by_t["won"]/by_t["pros"]*100, 0)
    by_t["CPA"]  = np.where(by_t["cost"]>0, by_t["cost"]/by_t["won"].replace(0,np.nan), 0)
    by_t = by_t.sort_values("E2E%", ascending=False).reset_index(drop=True)

    rect(s, Inches(9.4), Inches(1.12), Inches(3.7), Inches(5.95), LGRAY)
    hdrs = ["Team","Won","E2E%","CPA"]
    col_x = [Inches(9.5), Inches(10.5), Inches(11.25), Inches(12.1)]
    col_w = [Inches(0.9), Inches(0.65), Inches(0.75), Inches(0.9)]
    rect(s, Inches(9.4), Inches(1.12), Inches(3.7), Inches(0.4), DARK)
    for cx, cw, h in zip(col_x, col_w, hdrs):
        txt(s, h, cx, Inches(1.18), cw, Inches(0.3), size=9, bold=True, color=WHITE)
    for j, row in by_t.iterrows():
        y = Inches(1.52) + j * Inches(0.72)
        row_bg = LGRAY2 if j % 2 == 0 else LGRAY
        rect(s, Inches(9.4), y, Inches(3.7), Inches(0.68), row_bg)
        cpa_str = "—" if row["CPA"] == 0 else f"${row['CPA']:,.0f}"
        e2e_c = GREEN if row["E2E%"] >= 2 else (AMBER if row["E2E%"] >= 0.5 else RED_C)
        vals = [row["Marketing Team"], f"{row['won']:,.0f}", f"{row['E2E%']:.2f}%", cpa_str]
        for k, (cx, cw, v) in enumerate(zip(col_x, col_w, vals)):
            tc = e2e_c if k == 2 else DARK
            bold = k == 2
            txt(s, str(v), cx, y+Inches(0.18), cw, Inches(0.35), size=9, bold=bold, color=tc)

    insight_box(s, Inches(0.3), Inches(5.0), Inches(9.0), Inches(1.0),
                "Lifecycle leads with 4.0% E2E rate (highest of all teams) at ZERO cost. "
                "Paid has the most volume but only 0.13% E2E — 31x less efficient than lifecycle per prospect. "
                "Regional outperforms paid at the lowest CPA in the portfolio with far less investment.")

def slide_channel_conversion(prs, df):
    s = add_slide(prs)
    bg(s, WHITE)
    header_bar(s, "Task 1 — 3-Channel Conversion Rates",
               "Paid Search · Event · Regional  |  Impression → Click → Demo → SQO → Won")
    orange_bar(s)
    buf = chart_channel_conversion(df)
    img(s, buf, Inches(0.3), Inches(1.12), Inches(8.8), Inches(3.7))

    channels = ["paid_search", "event", "regional"]
    ch = (df[df["Marketing Program"].isin(channels)]
          .groupby("Marketing Program")
          .agg(imp=("# Impressions","sum"), clk=("# Clicks","sum"),
               demo=("# Demo Prospects","sum"), sqo=("# Sales Qualified Opportunities","sum"),
               won=("# Opportunities Won","sum"), cost=("Cost","sum")).reset_index())
    ch["CTR%"]      = np.where(ch["imp"]>0, ch["clk"]/ch["imp"]*100, np.nan)
    ch["Clk→Demo%"] = np.where(ch["clk"]>0, ch["demo"]/ch["clk"]*100, np.nan)
    ch["Demo→SQO%"] = np.where(ch["demo"]>0, ch["sqo"]/ch["demo"]*100, np.nan)
    ch["SQO→Won%"]  = np.where(ch["sqo"]>0, ch["won"]/ch["sqo"]*100, np.nan)
    ch["CPA"]       = np.where(ch["won"]>0, ch["cost"]/ch["won"], np.nan)
    ch_d = {row["Marketing Program"]: row for _, row in ch.iterrows()}

    rect(s, Inches(9.25), Inches(1.12), Inches(3.85), Inches(5.95), LGRAY)
    metrics = ["CTR%","Clk→Demo%","Demo→SQO%","SQO→Won%","CPA ($)","Total Spend"]
    col_names = ["Metric","Paid Search","Event","Regional"]
    col_x_t = [Inches(9.35), Inches(10.45), Inches(11.35), Inches(12.25)]
    rect(s, Inches(9.25), Inches(1.12), Inches(3.85), Inches(0.4), DARK)
    for cx, cn in zip(col_x_t, col_names):
        txt(s, cn, cx, Inches(1.17), Inches(0.9), Inches(0.35), size=8, bold=True, color=WHITE)

    def fmt_val(prog, metric):
        row = ch_d.get(prog, None)
        if row is None: return "—"
        if metric == "CTR%": return f"{row['CTR%']:.2f}%" if not np.isnan(row['CTR%']) else "n/a"
        if metric == "Clk→Demo%": return f"{row['Clk→Demo%']:.2f}%" if not np.isnan(row.get('Clk→Demo%', np.nan)) else "n/a"
        if metric == "Demo→SQO%": return f"{row['Demo→SQO%']:.1f}%" if not np.isnan(row['Demo→SQO%']) else "n/a"
        if metric == "SQO→Won%": return f"{row['SQO→Won%']:.1f}%" if not np.isnan(row['SQO→Won%']) else "n/a"
        if metric == "CPA ($)": return f"${row['CPA']:,.0f}" if not np.isnan(row['CPA']) else "—"
        if metric == "Total Spend": return f"${row['cost']/1e6:.1f}M"
        return "—"

    for j, metric in enumerate(metrics):
        y = Inches(1.52) + j * Inches(0.88)
        row_bg = LGRAY2 if j % 2 == 0 else LGRAY
        rect(s, Inches(9.25), y, Inches(3.85), Inches(0.84), row_bg)
        txt(s, metric, col_x_t[0], y+Inches(0.22), Inches(1.0), Inches(0.4),
            size=8.5, bold=True, color=DARK)
        for k, prog in enumerate(channels):
            v = fmt_val(prog, metric)
            txt(s, v, col_x_t[k+1], y+Inches(0.22), Inches(0.85), Inches(0.4),
                size=9, bold=(metric=="CPA ($)"), color=DARK)

    insight_box(s, Inches(0.3), Inches(5.0), Inches(8.8), Inches(1.0),
                "Regional is the standout: lowest CPA ($17K vs $47K paid search) with a strong 34% Demo→SQO rate, "
                "yet received only 1.2% of paid budget. Event drives quality demos (24% Demo→SQO) with zero impression cost. "
                "Paid search volume leader but efficiency gap vs regional is significant.")

def slide_t1_recommendations(prs):
    s = add_slide(prs)
    bg(s, WHITE)
    header_bar(s, "Task 1 — Recommendations: Improve Funnel Efficiency & Volume")
    orange_bar(s)
    recs = [
        ("1", "Reallocate paid_social → paid_brand_search",
         "Paid social CPA is $197K — 11× worse than brand search ($18K). "
         "Redirecting 30% of paid social budget (~$24M) to brand search could generate "
         "+1,300 additional Won per year at the same total spend."),
        ("2", "Scale regional demand gen (under-funded)",
         "Regional has the best CPA ($17K) and strong Demo→SQO (34%) "
         "yet only received $4.3M total (1.2% of paid budget). "
         "3× budget increase to $13M could yield ~750 additional Won."),
        ("3", "Invest in lifecycle / CRM automation",
         "Lifecycle holds the highest E2E rate (4.04%) at zero media cost. "
         "Investing in Braze sequence expansion and lead scoring would "
         "capture warm prospects that currently leak or convert via expensive paid."),
        ("4", "Protect Q4 spend; reduce Q1 scale",
         "Won peaks in Oct–Dec and May–Jun with strong CPA. "
         "Jan–Mar shows worst CPA ($50–65K). Shift Q1 budget to lifecycle and "
         "content to maintain pipeline at lower cost during slow conversion period."),
    ]
    for i, (num, title, body) in enumerate(recs):
        col = i % 2; row_i = i // 2
        x = Inches(0.35 + col * 6.55); y = Inches(1.2 + row_i * 2.9)
        rect(s, x, y, Inches(6.25), Inches(2.65), LGRAY)
        rect(s, x, y, Inches(0.55), Inches(2.65), ORANGE)
        txt(s, num, x+Inches(0.13), y+Inches(0.65), Inches(0.35), Inches(0.65),
            size=26, bold=True, color=WHITE)
        txt(s, title, x+Inches(0.7), y+Inches(0.12), Inches(5.35), Inches(0.5),
            size=12, bold=True, color=DARK)
        txt(s, body, x+Inches(0.7), y+Inches(0.68), Inches(5.35), Inches(1.75),
            size=10, color=MUTED)

def slide_seasonality_causes(prs):
    s = add_slide(prs)
    bg(s, WHITE)
    header_bar(s, "Task 1 — Seasonality: Patterns, Root Causes & Implications",
               "Why do metrics rise and fall? Business reasons behind the data")
    orange_bar(s)

    # Column definitions
    col_x  = [Inches(0.20), Inches(4.10), Inches(8.80)]
    col_w  = [Inches(3.75), Inches(4.55), Inches(4.25)]
    col_headers = ["Pattern Observed", "Probable Cause", "Implication for Marketing"]

    # Header row
    hdr_y = Inches(1.10); hdr_h = Inches(0.50)
    for x, w, label in zip(col_x, col_w, col_headers):
        rect(s, x, hdr_y, w, hdr_h, DARK)
        txt(s, label, x+Inches(0.12), hdr_y+Inches(0.08), w-Inches(0.2), hdr_h,
            size=10, bold=True, color=WHITE)

    rows = [
        ("Impressions spike in Jan 2024\n(3× above period average)",
         "Deliberate large-scale awareness campaign launched by the paid team — not organic seasonality",
         "Isolate campaign periods before attributing results to market behaviour"),
        ("Won drops every Jan–Mar\n(consistent both years)",
         "B2B buyers reset annual budgets in January; decision-making slows across the industry",
         "Avoid scaling paid CPA spend in Q1 — it's the most expensive period per Won"),
        ("Won peaks Oct–Dec\n(Q4 enterprise surge)",
         "Enterprise customers close deals before their own fiscal year-end to spend allocated budgets",
         "Front-load pipeline building in Aug–Sep to harvest demand when buyers are ready"),
        ("Won peaks again in May–Jun\n(second seasonal high)",
         "Deals entered in Nov–Dec mature after ~150-day avg sales cycle, closing in May–Jun",
         "Feed top-of-funnel aggressively in Nov so pipeline ripens for the May peak"),
        ("CTR declining across 2024–2025\n(paid_search & paid_social)",
         "Audience saturation and creative fatigue — same users seeing repetitive ads, ignoring them",
         "Rotate creatives, expand audience targeting, and shift budget to brand search where CTR is 6.7%"),
    ]

    row_h = Inches(1.03)
    for i, (pattern, cause, impl) in enumerate(rows):
        y = Inches(1.62) + i * row_h
        fill = LGRAY if i % 2 == 0 else WHITE
        for x, w in zip(col_x, col_w):
            rect(s, x, y, w, row_h - Inches(0.04), fill)

        # orange left accent on first column only
        rect(s, col_x[0], y, Inches(0.06), row_h - Inches(0.04), ORANGE)

        txt(s, pattern, col_x[0]+Inches(0.15), y+Inches(0.08),
            col_w[0]-Inches(0.2), row_h-Inches(0.12), size=9.5, bold=True, color=DARK)
        txt(s, cause,   col_x[1]+Inches(0.10), y+Inches(0.08),
            col_w[1]-Inches(0.15), row_h-Inches(0.12), size=9, color=MUTED)
        txt(s, impl,    col_x[2]+Inches(0.10), y+Inches(0.08),
            col_w[2]-Inches(0.15), row_h-Inches(0.12), size=9, color=NAVY)

def slide_paid_channels(prs, df):
    s = add_slide(prs)
    bg(s, WHITE)
    header_bar(s, "Task 2 — Paid Search vs Brand Search: Cost → Demos → SQOs → Won",
               "Two core paid channels analyzed across the full period")
    orange_bar(s)
    buf = chart_paid_search_brand(df)
    img(s, buf, Inches(0.3), Inches(1.12), Inches(12.7), Inches(5.8))

def slide_cpm_ctr(prs, df):
    s = add_slide(prs)
    bg(s, WHITE)
    header_bar(s, "Task 2 — CPM vs CTR Matrix: Paid Channel Efficiency",
               "Upper-left quadrant (low CPM, high CTR) = most efficient awareness spend")
    orange_bar(s)
    buf = chart_cpm_ctr(df)
    img(s, buf, Inches(0.3), Inches(1.12), Inches(8.8), Inches(4.5))

    paid = df[df["Marketing Team"]=="paid"]
    imp_df = (paid[paid["# Impressions"]>0]
              .groupby("Marketing Program")
              .agg(imp=("# Impressions","sum"), clk=("# Clicks","sum"),
                   cost=("Cost","sum"), won=("# Opportunities Won","sum")).reset_index())
    imp_df["CTR"] = imp_df["clk"]/imp_df["imp"]*100
    imp_df["CPM"] = imp_df["cost"]/imp_df["imp"]*1000
    imp_df["CPA"] = np.where(imp_df["won"]>0, imp_df["cost"]/imp_df["won"], np.nan)

    rect(s, Inches(9.25), Inches(1.12), Inches(3.85), Inches(5.95), LGRAY)
    txt(s, "Channel Efficiency Summary", Inches(9.4), Inches(1.22), Inches(3.6), Inches(0.35),
        size=11, bold=True, color=DARK)
    best_ctr = imp_df.loc[imp_df["CTR"].idxmax()]
    best_cpm = imp_df.loc[imp_df["CPM"].idxmin()]
    worst_cpa = imp_df.dropna(subset=["CPA"]).loc[imp_df.dropna(subset=["CPA"])["CPA"].idxmax()]
    best_cpa  = imp_df.dropna(subset=["CPA"]).loc[imp_df.dropna(subset=["CPA"])["CPA"].idxmin()]

    def _fmt_cpa(v):
        if v >= 1000: return f"${v:,.0f}"
        return f"${v:.2f}"

    callouts = [
        ("Highest CTR", f"{best_ctr['Marketing Program']}  {best_ctr['CTR']:.2f}%", GREEN),
        ("Lowest CPM",  f"{best_cpm['Marketing Program']}  ${best_cpm['CPM']:.2f}", GREEN),
        ("Best CPA",    f"{best_cpa['Marketing Program']}  {_fmt_cpa(best_cpa['CPA'])}", GREEN),
        ("Worst CPA",   f"{worst_cpa['Marketing Program']}  {_fmt_cpa(worst_cpa['CPA'])}", RED_C),
    ]
    for j, (label, val, c) in enumerate(callouts):
        y = Inches(1.65) + j * Inches(1.28)
        rect(s, Inches(9.35), y, Inches(3.6), Inches(1.1), WHITE)
        rect(s, Inches(9.35), y, Inches(0.06), Inches(1.1), c)
        txt(s, label, Inches(9.5), y+Inches(0.08), Inches(3.3), Inches(0.3),
            size=9, color=MUTED)
        txt(s, val, Inches(9.5), y+Inches(0.4), Inches(3.3), Inches(0.55),
            size=12, bold=True, color=DARK)

    insight_box(s, Inches(0.3), Inches(5.35), Inches(8.8), Inches(1.0),
                "Brand search dominates on CTR (6.68%) — direct pull of branded intent. "
                "Paid awareness has the worst CPM efficiency (high cost, minimal conversions). "
                "Paid social shows 0.47% CTR despite being the highest-spend program — clear candidate for budget cut.")

def slide_brand_channels(prs, df):
    s = add_slide(prs)
    bg(s, WHITE)
    header_bar(s, "Task 2 — Brand Channels: Content, Website & Lifecycle Growth",
               "Zero-cost channels building long-term equity")
    orange_bar(s)
    buf = chart_brand_channels(df)
    img(s, buf, Inches(0.15), Inches(1.12), Inches(13.0), Inches(4.2))

    cyc_data = df[df["Cycle"].isin(["Y1 (Jul23-Jun24)","Y2 (Jul24-Jun25)"])].copy()
    grp = (cyc_data.groupby(["Cycle","Marketing Team"])
           .agg(won=("# Opportunities Won","sum"), pros=("# Prospects","sum")).reset_index())

    def yoy(team, metric):
        y1 = grp[(grp["Cycle"]=="Y1 (Jul23-Jun24)") & (grp["Marketing Team"]==team)][metric].values
        y2 = grp[(grp["Cycle"]=="Y2 (Jul24-Jun25)") & (grp["Marketing Team"]==team)][metric].values
        v1 = y1[0] if len(y1) else 0
        v2 = y2[0] if len(y2) else 0
        pct = (v2/v1-1)*100 if v1 > 0 else 0
        return v1, v2, pct

    paid_cpa = df[df["Cost"]>0]["Cost"].sum() / max(df[df["Cost"]>0]["# Opportunities Won"].sum(), 1)

    w1c, w2c, pc  = yoy("content",   "won")
    w1w, w2w, pw  = yoy("website",   "won")
    w1l, w2l, pl  = yoy("lifecycle", "won")
    p1c, p2c, ppc = yoy("content",   "pros")
    impl_val = (w2c + w2w + w2l) * paid_cpa

    kpis = [
        ("Content Won\nY1→Y2", f"{w1c:.0f}→{w2c:.0f}", f"+{pc:.0f}%", ORANGE),
        ("Website Won\nY1→Y2", f"{w1w:.0f}→{w2w:.0f}", f"{pw:+.0f}%", BLUE_C),
        ("Lifecycle Won\nY1→Y2", f"{w1l:.0f}→{w2l:.0f}", f"+{pl:.0f}%", PURPLE),
        ("Implied Value\n@ paid CPA", f"${impl_val/1e6:.1f}M", "zero spend", GREEN),
    ]
    for i, (label, val, delta, c) in enumerate(kpis):
        x = Inches(0.3 + i * 3.25)
        rect(s, x, Inches(4.85), Inches(3.1), Inches(2.0), LGRAY)
        rect(s, x, Inches(4.85), Inches(3.1), Inches(0.07), c)
        txt(s, label, x+Inches(0.2), Inches(4.97), Inches(2.8), Inches(0.5),
            size=9, color=MUTED)
        txt(s, val, x+Inches(0.2), Inches(5.5), Inches(2.8), Inches(0.55),
            size=18, bold=True, color=DARK)
        txt(s, delta, x+Inches(0.2), Inches(6.1), Inches(2.8), Inches(0.4),
            size=12, bold=True, color=c)

def slide_brand_measurement(prs):
    s = add_slide(prs)
    bg(s, WHITE)
    header_bar(s, "Task 2 — How to Measure Brand Budget Impact",
               "If Deel increases brand budget next quarter — specific measurement framework")
    orange_bar(s)
    metrics = [
        ("01", "Branded Search Volume",
         "Track Google Search Console impressions for 'Deel' / 'Deel HR' weekly. "
         "15%+ lift in branded impressions within 60 days = awareness signal. "
         "Baseline: current monthly avg impressions by keyword."),
        ("02", "Direct Traffic to Site",
         "Monitor GA4 direct sessions. Brand spend should drive recall. "
         "Target: direct traffic growing ≥10% vs prior quarter. "
         "Split by device and market to isolate brand vs technical effects."),
        ("03", "Brand Search CPA Trend",
         "Increased share-of-voice → lower brand search CPA. "
         "Current baseline: $18.7K. Target: below $15K within 8 weeks. "
         "Track weekly in the dashboard to catch regression early."),
        ("04", "Content MQL Rate",
         "If brand investment improves inbound quality, content MQL rate should rise. "
         "Current Y2 baseline: ~1.4% E2E. Track monthly MQL% vs 6-month rolling avg. "
         "Leading indicator before Won changes appear."),
        ("05", "Pipeline Influenced %",
         "Tag brand touchpoints (content views, direct visits, brand search) in CRM. "
         "Measure % of closed-won deals with ≥1 brand touchpoint in the journey. "
         "Target: ≥40% of Won influenced within 90 days."),
        ("06", "Sales Cycle Length",
         "Stronger brand recognition shortens sales cycles. "
         "Track avg days from Demo Booked → Won as proxy. "
         "Baseline: measure current avg, track weekly after brand spend increase."),
    ]
    for i, (num, title, body) in enumerate(metrics):
        col = i % 2; row_i = i // 2
        x = Inches(0.35 + col * 6.55); y = Inches(1.15 + row_i * 1.98)
        rect(s, x, y, Inches(6.25), Inches(1.82), LGRAY)
        txt(s, num, x+Inches(0.15), y+Inches(0.15), Inches(0.8), Inches(0.55),
            size=20, bold=True, color=ORANGE)
        txt(s, title, x+Inches(0.9), y+Inches(0.1), Inches(5.1), Inches(0.45),
            size=12, bold=True, color=DARK)
        txt(s, body, x+Inches(0.15), y+Inches(0.65), Inches(5.95), Inches(1.0),
            size=9.5, color=MUTED)

def slide_reallocation(prs, df):
    s = add_slide(prs)
    bg(s, WHITE)
    header_bar(s, "Task 2 — Paid Budget Reallocation Recommendation",
               "Where to scale, reduce, and cut — based on CPA efficiency tiers")
    orange_bar(s)
    buf = chart_reallocation(df)
    img(s, buf, Inches(0.3), Inches(1.12), Inches(8.8), Inches(4.2))

    # ── dynamic CPA labels from Cost column ─────────────────────────────────
    _pp = (df[df["Marketing Team"]=="paid"]
           .groupby("Marketing Program")
           .agg(cost=("Cost","sum"), won=("# Opportunities Won","sum")).reset_index())
    _pp["CPA"] = np.where(_pp["won"]>0, _pp["cost"]/_pp["won"], np.nan)
    def _cpa(prog, fallback="n/a"):
        r = _pp[_pp["Marketing Program"]==prog]["CPA"]
        return f"${r.values[0]:,.0f}" if len(r) and not np.isnan(r.values[0]) else fallback
    def _mult(p1, p2):
        c1 = _pp.loc[_pp["Marketing Program"]==p1,"CPA"]
        c2 = _pp.loc[_pp["Marketing Program"]==p2,"CPA"]
        if len(c1) and len(c2) and not np.isnan(c1.values[0]) and c2.values[0]>0:
            return f"{c1.values[0]/c2.values[0]:.0f}×"
        return "many×"
    # regional may sit outside "paid" team
    _reg_sub = df[df["Marketing Program"]=="regional"]
    _reg_cost = _reg_sub["Cost"].sum()
    _reg_won  = _reg_sub["# Opportunities Won"].sum()
    _reg_cpa  = f"${_reg_cost/_reg_won:,.0f}" if _reg_won > 0 else "$17"
    _freed = df[df["Marketing Program"].isin(["paid_social","paid_awareness"])]["Cost"].sum() * 0.5

    decisions = [
        ("SCALE UP", [
            f"paid_brand_search — CPA {_cpa('paid_brand_search')}, CTR 6.68%",
            f"regional — CPA {_reg_cpa}, Demo→SQO 34%, severely under-funded",
        ], GREEN),
        ("REDUCE / OPTIMIZE", [
            f"paid_search — CPA {_cpa('paid_search')} (high vol.), watch CPA trend",
            f"listicle/leadbuy — CPA {_cpa('listicle_leadbuy_listing')}, very low CTR (0.03%)",
        ], AMBER),
        ("CUT", [
            f"paid_social — CPA {_cpa('paid_social')}  ({_mult('paid_social','paid_brand_search')} brand search)",
            f"paid_awareness — CPA {_cpa('paid_awareness')}  (only 3 Won from entire program)",
            f"paid_sponsorship — CPA {_cpa('paid_sponsorship')}, minimal Won volume (27 total)",
        ], RED_C),
    ]
    y_start = Inches(1.15)
    for label, items, color in decisions:
        rect(s, Inches(9.25), y_start, Inches(3.85), Inches(0.38), color)
        txt(s, label, Inches(9.4), y_start+Inches(0.06), Inches(3.5), Inches(0.3),
            size=12, bold=True, color=WHITE)
        for j, item in enumerate(items):
            txt(s, f"→  {item}", Inches(9.4), y_start+Inches(0.42+j*0.44),
                Inches(3.65), Inches(0.4), size=9.5, color=DARK)
        y_start += Inches(0.38 + len(items)*0.44 + 0.28)

    insight_box(s, Inches(0.3), Inches(5.45), Inches(8.8), Inches(1.1),
                f"Reallocating 50% of paid_social + paid_awareness (${_freed/1e6:.0f}M freed) "
                f"to brand_search and regional at historical CPAs is projected to generate "
                f"+23% additional Won Opps per year with zero budget increase.")

def slide_budget_impact(prs, df):
    s = add_slide(prs)
    bg(s, WHITE)
    header_bar(s, "Task 2 — Budget Reallocation: Projected Impact",
               "Scenario: cut high-CPA programs, reinvest in low-CPA programs")
    orange_bar(s)
    buf = chart_budget_impact(df)
    img(s, buf, Inches(0.3), Inches(1.12), Inches(12.7), Inches(4.2))

    paid = (df[df["Marketing Team"]=="paid"].groupby("Marketing Program")
            .agg(cost=("Cost","sum"), won=("# Opportunities Won","sum")).reset_index())
    paid["CPA"] = np.where(paid["won"]>0, paid["cost"]/paid["won"], np.nan)
    paid = paid.dropna(subset=["CPA"]).sort_values("CPA").reset_index(drop=True)
    n = len(paid)
    n_low  = max(n // 3, 1)
    n_high = max(n // 3, 1)
    low_cpa  = paid.iloc[:n_low]
    high_cpa = paid.iloc[n - n_high:]
    mid_cpa  = paid.iloc[n_low : n - n_high]
    low_thresh  = low_cpa["CPA"].max()
    high_thresh = high_cpa["CPA"].min()
    cut_val  = high_cpa["cost"].sum() * 0.6
    add_won  = (cut_val / low_cpa["CPA"].mean()) if low_cpa["CPA"].mean() > 0 else 0
    high_mean = high_cpa["CPA"].mean() if len(high_cpa) > 0 else 1
    won_curr = paid["won"].sum()
    won_opt  = (low_cpa["won"].sum() + add_won
                + mid_cpa["won"].sum()
                + high_cpa["cost"].sum() * 0.4 / max(high_mean, 0.01))

    assumptions = [
        f"Cut 60% from highest-CPA programs (CPA > ${high_thresh/1e3:.0f}K) — saves ~${cut_val/1e6:.0f}M",
        f"Reinvest in low-CPA programs (CPA ≤ ${low_thresh/1e3:.0f}K) at historical efficiency",
        "Mid-tier programs keep current allocation",
        "CPA assumed stable — actual may improve with better audience targeting",
    ]
    rect(s, Inches(0.3), Inches(5.45), Inches(12.7), Inches(1.6), LGRAY)
    txt(s, "Scenario Assumptions", Inches(0.5), Inches(5.55), Inches(12.3), Inches(0.32),
        size=11, bold=True, color=DARK)
    for j, a in enumerate(assumptions):
        col = j % 2; row_i = j // 2
        x = Inches(0.5 + col * 6.5); y = Inches(5.92) + row_i * Inches(0.42)
        rect(s, x-Inches(0.1), y+Inches(0.06), Inches(0.06), Inches(0.2), ORANGE)
        txt(s, a, x, y, Inches(6.2), Inches(0.38), size=9.5, color=DARK)

def slide_scenario_paid(prs):
    s = add_slide(prs)
    bg(s, WHITE)
    header_bar(s, "Task 3 — Executive Summary: Paid Marketing Director",
               "ROI-focused takeaways and recommended actions")
    orange_bar(s)
    rect(s, 0, Inches(1.05), Inches(0.12), H-Inches(1.11), ORANGE)

    sections_data = [
        ("Executive Summary (max 150 words)",
         "Our paid portfolio generated 8,942 Won at a blended CPA of $40K, "
         "but efficiency varies 112× across programs. Paid brand search ($19K CPA, 6.7% CTR) "
         "outperforms paid social ($198K CPA) and paid awareness ($2.1M CPA) by an enormous margin. "
         "Brand search CPA has been rising through 2025 — a signal of growing audience saturation "
         "at current spend levels. Regional ($17K CPA) is our most efficient demand "
         "gen channel yet receives only 1.2% of paid budget. Without reallocation, "
         "we are overpaying for low-quality conversions at scale.", DARK),
        ("Recommended Action",
         "Immediate: (1) Cut paid_social and paid_awareness by 50% (~$44M freed). "
         "(2) Redirect freed budget to paid_brand_search and regional. At historical CPAs, this generates "
         "+23% additional Won per year at zero budget increase. "
         "(3) Monitor paid_search monthly — CPA rising in 2025 indicates diminishing returns at current scale.", ORANGE),
        ("Measuring Success (next 1–2 months)",
         "Week 1–2: Pause paid_social/awareness. Track weekly Won — expect initial dip as pipeline adjusts. "
         "Week 3–4: Monitor brand_search CPA daily; target staying below $25K. "
         "Month 2: Compare total Won and blended CPA vs prior 30-day baseline. "
         "Success = same or higher Won count at ≥15% lower blended CPA.", GREEN),
    ]
    sec_colors = [DARK, ORANGE, GREEN]
    y = Inches(1.15)
    for (title, body, _), sec_c in zip(sections_data, sec_colors):
        rect(s, Inches(0.25), y, Inches(12.8), Inches(0.38), sec_c)
        txt(s, title, Inches(0.4), y+Inches(0.05), Inches(12.4), Inches(0.3),
            size=11, bold=True, color=WHITE)
        box_h = Inches(1.22)
        rect(s, Inches(0.25), y+Inches(0.38), Inches(12.8), box_h, LGRAY)
        txt(s, body, Inches(0.4), y+Inches(0.48), Inches(12.5), box_h-Inches(0.12),
            size=11, color=DARK)
        y += Inches(0.38) + box_h + Inches(0.15)

def slide_scenario_head(prs):
    s = add_slide(prs)
    bg(s, WHITE)
    header_bar(s, "Task 3 — Executive Summary: Head of Marketing",
               "Brand and long-term growth takeaways and recommended actions")
    orange_bar(s)
    rect(s, 0, Inches(1.05), Inches(0.12), H-Inches(1.11), BLUE_C)

    sections_data = [
        ("Executive Summary (max 150 words)",
         "Zero-cost organic channels are quietly outperforming expectations. "
         "Lifecycle (CRM/email) grew +132% YoY in Won at zero media spend — "
         "holding a 4.04% E2E conversion rate, the highest of any channel in the portfolio. "
         "Content and website together contribute ~580 Won annually — at the blended paid CPA ($40K), "
         "that represents ~$23M in acquisition value delivered at zero media cost. "
         "Brand search's 6.68% CTR vs performance channels (0.47–3.6%) "
         "proves that brand investment creates pull that measurably reduces acquisition cost over time. "
         "Content prospect volume grew 146% YoY — a leading indicator of future pipeline health.", DARK),
        ("Recommended Action",
         "Allocate 10–15% of paid budget to lifecycle automation (lead scoring, email sequencing) "
         "to capitalize on the 4.04% E2E rate at scale — the highest-ROI investment available in the portfolio. "
         "Increase content production focused on SEO and thought leadership to sustain 32K+ prospect/month volume. "
         "Protect brand search budget to defend the low-CPA, high-CTR branded channel from erosion.", BLUE_C),
        ("Measuring Success (next 1–2 months)",
         "Month 1: Baseline lifecycle Won/month (current ~40/mo). After automation, target +25% in 60 days. "
         "Track branded search impression share in Google Search Console — target +10% lift. "
         "Monitor content E2E rate monthly vs 1.41% Y2 baseline. "
         "Long-term KPI: % of closed-won deals with ≥1 content or direct-brand touchpoint.", BLUE_C),
    ]
    sec_colors = [DARK, BLUE_C, BLUE_C]
    y = Inches(1.15)
    for (title, body, _), sec_c in zip(sections_data, sec_colors):
        rect(s, Inches(0.25), y, Inches(12.8), Inches(0.38), sec_c)
        txt(s, title, Inches(0.4), y+Inches(0.05), Inches(12.4), Inches(0.3),
            size=11, bold=True, color=WHITE)
        box_h = Inches(1.22)
        rect(s, Inches(0.25), y+Inches(0.38), Inches(12.8), box_h, LGRAY)
        txt(s, body, Inches(0.4), y+Inches(0.48), Inches(12.5), box_h-Inches(0.12),
            size=11, color=DARK)
        y += Inches(0.38) + box_h + Inches(0.15)

def slide_action_plan(prs):
    s = add_slide(prs)
    bg(s, WHITE)
    header_bar(s, "Executive Action Plan — 90-Day Roadmap",
               "Priority actions across both stakeholder perspectives")
    orange_bar(s)

    phases = [
        ("Immediate (Week 1–2)", ORANGE, [
            ("Pause paid_social + paid_awareness campaigns",
             "Free ~$44M from programs above the $50K CPA benchmark. Zero impact on pipeline within 2 weeks."),
            ("Activate regional expansion budget",
             "Double regional allocation from $4.3M → $8.6M. Best CPA in portfolio at $17K."),
        ]),
        ("Short-term (Month 1–2)", BLUE_C, [
            ("Ramp paid_brand_search",
             "Redirect freed budget from cut programs. Monitor CTR + CPA weekly. Target CPA <$25K."),
            ("Launch lifecycle automation build",
             "Engage Braze team to expand sequences. Prioritize warm MQPs with demo intent."),
        ]),
        ("Strategic (Month 2–3)", GREEN, [
            ("Set quarterly CPA targets by program",
             "Establish $50K ceiling for scale-up programs; alert if CPA exceeds target by >20%."),
            ("Establish brand measurement framework",
             "Launch brand search tracking in GSC, pipeline influence tagging in CRM."),
        ]),
    ]
    col_w = Inches(4.2); col_gap = Inches(0.17)
    for i, (phase_title, c, actions) in enumerate(phases):
        x = Inches(0.3) + i*(col_w+col_gap)
        rect(s, x, Inches(1.12), col_w, Inches(0.45), c)
        txt(s, phase_title, x+Inches(0.15), Inches(1.17), col_w-Inches(0.2), Inches(0.35),
            size=12, bold=True, color=WHITE)
        for j, (action_title, action_body) in enumerate(actions):
            y = Inches(1.62) + j*Inches(2.6)
            rect(s, x, y, col_w, Inches(2.45), LGRAY)
            rect(s, x, y, Inches(0.06), Inches(2.45), c)
            txt(s, action_title, x+Inches(0.15), y+Inches(0.1), col_w-Inches(0.25),
                Inches(0.45), size=11, bold=True, color=DARK)
            txt(s, action_body, x+Inches(0.15), y+Inches(0.58), col_w-Inches(0.25),
                Inches(1.7), size=10, color=MUTED)

def slide_appendix(prs):
    s = add_slide(prs)
    bg(s, WHITE)
    header_bar(s, "Appendix — Methodology, Tools & Data Notes")
    orange_bar(s)
    items = [
        ("Dataset", "11,952 rows | Jul 2023 – Jun 2025 | 7 marketing teams | 17 programs | "
                    "Semicolon-delimited CSV | Cost strings cleaned with regex (EU & US formats handled)"),
        ("Tools Used", "Python 3.10 · Pandas (wrangling) · NumPy (metrics) · "
                       "Plotly + Streamlit (interactive dashboard, 6 tabs) · "
                       "Matplotlib + python-pptx (this presentation)"),
        ("Attribution Note", "Opportunity Won values are decimals due to multi-touch attribution model. "
                             "Values were summed (not rounded) throughout. Decimal Wons are expected."),
        ("Fiscal Year", "Analysis uses Jul–Jun cycles: Y1 = Jul 2023–Jun 2024, Y2 = Jul 2024–Jun 2025."),
        ("CPA Calculation", "CPA = Total Cost ÷ Total Won per program/team. "
                            "Zero-cost channels (lifecycle, content, website, social) excluded from CPA."),
        ("Implied Value", "Zero-cost channel implied value = Won × Avg Paid CPA. "
                          "Represents the cost Deel would have paid to generate the same Won via paid channels."),
        ("MQP Data Note", "MQP column was 0 for most of H2 2023 — data collection likely started in early 2024. "
                          "E2E rates for 2023 use Prospects→Won only to avoid division-by-zero distortion."),
        ("Interactive Dashboard", "A full Streamlit app (app.py) is included: Executive View, Full Funnel, "
                                   "Paid Channels, Zero-Cost Channels, Trends, Budget Simulator."),
    ]
    y = Inches(1.2)
    for label, body in items:
        rect(s, Inches(0.3), y, Inches(0.07), Inches(0.5), ORANGE)
        txt(s, label, Inches(0.5), y, Inches(12.5), Inches(0.3),
            size=11, bold=True, color=DARK)
        txt(s, body, Inches(0.5), y+Inches(0.3), Inches(12.5), Inches(0.38),
            size=9.5, color=MUTED)
        y += Inches(0.72)

def slide_t3_recommendations(prs, df):
    s = add_slide(prs)
    bg(s, WHITE)
    header_bar(s, "Task 3 — Specific Recommendations by Stakeholder",
               "Clear, data-backed actions each leader should take immediately")
    orange_bar(s)

    # ── compute live numbers for labels ──────────────────────────────────────
    paid = (df[df["Marketing Team"]=="paid"]
            .groupby("Marketing Program")
            .agg(cost=("Cost","sum"), won=("# Opportunities Won","sum")).reset_index())
    paid["CPA"] = np.where(paid["won"]>0, paid["cost"]/paid["won"], np.nan)

    social_cpa = paid.loc[paid["Marketing Program"]=="paid_social","CPA"].values
    social_cpa = f"${social_cpa[0]:,.0f}" if len(social_cpa) else "$197"
    brand_cpa  = paid.loc[paid["Marketing Program"]=="paid_brand_search","CPA"].values
    brand_cpa  = f"${brand_cpa[0]:,.0f}" if len(brand_cpa) else "$18"
    reg_cpa    = paid.loc[paid["Marketing Program"]=="regional","CPA"].values
    reg_cpa    = f"${reg_cpa[0]:,.0f}" if len(reg_cpa) else "$17"

    lc = df[df["Marketing Team"]=="lifecycle"]
    lc_e2e = lc["# Opportunities Won"].sum() / max(lc["# Prospects"].sum(),1) * 100

    # ── panel definitions ─────────────────────────────────────────────────────
    panels = [
        {
            "title":    "Paid Marketing Director",
            "tag":      "ROI  ·  SHORT-TERM PERFORMANCE",
            "color":    ORANGE,
            "steps": [
                ("1", "Cut paid_social & paid_awareness — this week",
                 f"CPA of {social_cpa} (social) makes these the worst-ROI programs in the portfolio. "
                 "Pause both immediately. The pipeline impact is negligible within 30 days "
                 "because these channels convert at <0.5% CTR and sit at the far top of funnel."),
                ("2", f"Redirect freed budget → brand_search & regional",
                 f"paid_brand_search CPA is {brand_cpa} (CTR 6.7%) and regional CPA is {reg_cpa} "
                 "(Demo→SQO 34%). Reallocating 50% of cut budget to these two programs "
                 "is projected to generate +23% Won at the same total spend."),
                ("3", "Set a hard CPA ceiling of $50K per program",
                 "Any program that exceeds $50K CPA for 2 consecutive months is automatically "
                 "paused and reviewed. Install a weekly CPA alert in the dashboard. "
                 "This creates a self-correcting budget governance rule going forward."),
            ],
            "kpis": [
                ("Target CPA",    "≤ $22K",       ORANGE),
                ("Won uplift",    "+23%",          GREEN),
                ("Timeline",      "2 weeks",       BLUE_C),
            ],
        },
        {
            "title":    "Head of Marketing",
            "tag":      "BRAND  ·  LONG-TERM EQUITY",
            "color":    BLUE_C,
            "steps": [
                ("1", "Commission lifecycle automation — highest-ROI investment available",
                 f"Lifecycle holds a {lc_e2e:.1f}% E2E rate — the best of any channel — at zero media cost. "
                 "Investing in Braze sequence expansion and lead scoring for MQPs "
                 "who haven't booked a demo within 7 days is the single highest-ROI action available."),
                ("2", "Scale content production with SEO-first briefs",
                 "Content prospects grew +146% YoY, feeding the lifecycle pipeline that converts at 4×+ "
                 "the paid average. Increase content output by 20% with keyword-targeted briefs. "
                 "This compounds organically — every piece published reduces paid dependency."),
                ("3", "Launch the 6-KPI brand measurement dashboard",
                 "Use the framework defined in slide 16 (branded search vol., direct traffic, "
                 "brand search CPA, content MQL rate, pipeline influence %, sales cycle length). "
                 "Establish baselines in week 1 so Q3 brand investment is fully accountable to the board."),
            ],
            "kpis": [
                ("Lifecycle Won",  "+25% in 60d",  BLUE_C),
                ("Brand proof",    "90 days",       GREEN),
                ("Media cost",     "$0",            GREEN),
            ],
        },
    ]

    for col_i, p in enumerate(panels):
        px = Inches(0.25 + col_i * 6.6)
        pw = Inches(6.45)
        ph = Inches(5.95)
        py = Inches(1.12)

        # panel background
        rect(s, px, py, pw, ph, LGRAY)
        rect(s, px, py, pw, Inches(0.55), p["color"])
        txt(s, p["title"], px+Inches(0.15), py+Inches(0.06),
            pw-Inches(0.25), Inches(0.35), size=14, bold=True, color=WHITE)
        txt(s, p["tag"],   px+Inches(0.15), py+Inches(0.38),
            pw-Inches(0.25), Inches(0.2),  size=8,  bold=True, color=WHITE)

        # action steps
        for j, (num, title, body) in enumerate(p["steps"]):
            sy = py + Inches(0.65) + j * Inches(1.69)
            rect(s, px+Inches(0.15), sy, pw-Inches(0.3), Inches(1.6), WHITE)
            rect(s, px+Inches(0.15), sy, Inches(0.06), Inches(1.6), p["color"])
            # number badge
            rect(s, px+Inches(0.22), sy+Inches(0.1), Inches(0.38), Inches(0.38), p["color"])
            txt(s, num, px+Inches(0.22), sy+Inches(0.1), Inches(0.38), Inches(0.38),
                size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            txt(s, title, px+Inches(0.68), sy+Inches(0.08),
                pw-Inches(1.0), Inches(0.38), size=10, bold=True, color=DARK)
            txt(s, body,  px+Inches(0.68), sy+Inches(0.48),
                pw-Inches(1.0), Inches(1.05), size=8.5, color=MUTED)

        # KPI strip at bottom of panel
        kpy = py + ph - Inches(0.65)
        rect(s, px, kpy, pw, Inches(0.65), p["color"])
        for k, (klabel, kval, kc) in enumerate(p["kpis"]):
            kx = px + Inches(0.2) + k * Inches(2.1)
            txt(s, klabel, kx, kpy+Inches(0.04), Inches(2.0), Inches(0.25),
                size=7.5, color=WHITE)
            txt(s, kval,   kx, kpy+Inches(0.28), Inches(2.0), Inches(0.32),
                size=14, bold=True, color=WHITE)


def slide_thank_you(prs):
    photo = os.path.join(os.path.dirname(os.path.abspath(__file__)), "slide24_photo.png")
    s = add_slide(prs)
    bg(s, DARK)

    # Person photo as full background
    if os.path.exists(photo):
        with open(photo, "rb") as f:
            img(s, f, Inches(0), Inches(0), W, H)

    # Dark panel covers right 55% — hides Wise branding in the photo
    rect(s, Inches(5.85), Inches(0), Inches(7.48), H, DARK)

    # Orange vertical divider line
    rect(s, Inches(5.85), Inches(0), Inches(0.07), H, ORANGE)

    # Cover Wise card in person's hand with a Deel-branded card shape
    rect(s, Inches(1.55), Inches(3.70), Inches(1.62), Inches(1.02),
         RGBColor(0xFF, 0x5C, 0x35))
    rect(s, Inches(1.55), Inches(3.70), Inches(1.62), Inches(0.07),
         RGBColor(0xFF, 0x8A, 0x65))
    # chip detail
    rect(s, Inches(1.68), Inches(3.85), Inches(0.32), Inches(0.22),
         RGBColor(0xFF, 0xC2, 0x9E))
    txt(s, "deel", Inches(1.60), Inches(4.30), Inches(1.45), Inches(0.32),
        size=13, bold=True, color=WHITE)

    # ── Right side Deel branding ──────────────────────────────────────────────
    # Deel logo
    txt(s, "deel", Inches(6.3), Inches(0.45), Inches(3.5), Inches(0.72),
        size=42, bold=True, color=ORANGE)

    # "Thank you!" heading
    txt(s, "Thank you!", Inches(6.3), Inches(1.38), Inches(6.7), Inches(1.3),
        size=54, bold=True, color=WHITE)

    # Orange underline accent
    rect(s, Inches(6.3), Inches(2.78), Inches(1.4), Inches(0.07), ORANGE)

    # Body text
    txt(s, "Thank you for your time and attention.\nI appreciate your interest and engagement.",
        Inches(6.3), Inches(3.0), Inches(6.7), Inches(0.9),
        size=13, color=RGBColor(0xBB, 0xBB, 0xCC))

    txt(s, "I'm happy to take any questions\nor hear your thoughts.",
        Inches(6.3), Inches(4.05), Inches(6.7), Inches(0.75),
        size=13, color=RGBColor(0xBB, 0xBB, 0xCC))

    # Tagline with orange accent bar
    rect(s, Inches(6.3), Inches(5.3), Inches(0.07), Inches(0.65), ORANGE)
    txt(s, "Simplifying global HR,\npayroll & compliance for teams everywhere.",
        Inches(6.5), Inches(5.3), Inches(6.5), Inches(0.72),
        size=10, color=MUTED)

    # deel.com
    txt(s, "deel.com", Inches(11.8), Inches(7.1), Inches(1.4), Inches(0.32),
        size=10, color=MUTED)

    # Bottom orange bar
    orange_bar(s)


# ── MAIN ──────────────────────────────────────────────────────────────────────
def main():
    print("Loading data...")
    df = load()
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    print("Building slides...")
    slide_cover(prs)                          # 1  — Cover
    slide_audience(prs)                       # 2  — Who you're presenting to
    slide_global_kpis(prs, df)               # 3  — Global KPIs + funnel waterfall

    # ── Task 1
    slide_task_divider(prs, 1,
        "Analyzing Sales Funnel Data",
        "Patterns over time · Seasonality · Channel conversion rates · Recommendations")
    slide_top_funnel(prs, df)                # 5  — Impressions & Clicks
    slide_bottom_funnel(prs, df)             # 6  — Demos, SQOs, Won over time
    slide_seasonality(prs, df)               # 7  — YoY seasonality
    slide_team_funnel(prs, df)               # 8  — E2E & CPA by team
    slide_channel_conversion(prs, df)        # 9  — 3-channel conversion
    slide_t1_recommendations(prs)            # 10 — Task 1 recommendations
    slide_seasonality_causes(prs)            # 11 — Seasonality root causes

    # ── Task 2
    slide_task_divider(prs, 2,
        "Paid & Brand Marketing Channel Analysis",
        "Cost vs outcomes · CPM/CTR efficiency · Brand growth · Budget reallocation")
    slide_paid_channels(prs, df)             # 13 — Paid search vs brand search
    slide_cpm_ctr(prs, df)                   # 14 — CPM vs CTR matrix
    slide_brand_channels(prs, df)            # 15 — Content / website / lifecycle growth
    slide_brand_measurement(prs)             # 16 — How to measure brand impact
    slide_reallocation(prs, df)              # 17 — Scale / reduce / cut
    slide_budget_impact(prs, df)             # 18 — Budget reallocation projection

    # ── Task 3
    slide_task_divider(prs, 3,
        "Scenario Analysis",
        "Stakeholder-specific summaries · Recommendations · Success metrics")
    slide_scenario_paid(prs)                 # 20 — Paid Marketing Director
    slide_scenario_head(prs)                 # 21 — Head of Marketing
    slide_action_plan(prs)                   # 22 — 90-day action roadmap
    slide_t3_recommendations(prs, df)        # 23 — Specific recommendations per stakeholder

    slide_appendix(prs)                      # 24 — Methodology & tools
    slide_thank_you(prs)                     # 25 — Thank you (Deel branded)

    out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                       "Deel_Marketing_Analytics_v2.pptx")
    prs.save(out)
    print(f"\nDone! Saved to: {out}")
    print(f"Slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
