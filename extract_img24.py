"""Extract the picture from slide 24 so generate_ppt.py can use it."""
from pptx import Presentation
import os

prs = Presentation(r"C:\Users\gabri\Downloads\deel\Deel_Marketing_Analytics_v2.pptx")
s = prs.slides[23]

for sh in s.shapes:
    if sh.shape_type == 13:   # PICTURE
        img_bytes = sh.image.blob
        ext = sh.image.ext
        out = rf"C:\Users\gabri\Downloads\deel\slide24_photo.{ext}"
        with open(out, "wb") as f:
            f.write(img_bytes)
        print(f"Saved: {out}  ({len(img_bytes)//1024} KB, ext={ext})")
        break
